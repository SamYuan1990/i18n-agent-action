import logging
import os

from pdfminer.high_level import extract_text
from pptx import Presentation


def is_pdf_file(filepath):
    """
    检查文件是否以.pdf结尾（不区分大小写）

    参数:
    filepath (str): 文件的绝对路径

    返回:
    bool: 如果是PDF文件返回True，否则返回False
    """
    # 使用os.path.splitext获取文件扩展名并转换为小写进行比较
    return os.path.splitext(filepath)[1].lower() == ".pdf"


def is_pptx_file(filepath):
    """
    检查文件是否以.pdf结尾（不区分大小写）

    参数:
    filepath (str): 文件的绝对路径

    返回:
    bool: 如果是PDF文件返回True，否则返回False
    """
    # 使用os.path.splitext获取文件扩展名并转换为小写进行比较
    return os.path.splitext(filepath)[1].lower() == ".pptx"


def pptx_to_markdown_advanced(pptx_path):
    prs = Presentation(pptx_path)
    markdown_lines = []

    for i, slide in enumerate(prs.slides):
        markdown_lines.append(f"# Slide {i+1}\n\n")

        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
            markdown_lines.append(f"## {title}\n\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    if any(char in text for char in ["•", "-", "*", "→"]):
                        for line in text.split("\n"):
                            if line.strip():
                                cleaned = line.strip().lstrip("•-*→ ")
                                markdown_lines.append(f"* {cleaned}\n")
                        markdown_lines.append("\n")
                    else:
                        markdown_lines.append(f"{text}\n\n")

            if shape.has_table:
                table = shape.table
                headers = [cell.text for cell in table.rows[0].cells]
                markdown_lines.append(f"| {' | '.join(headers)} |\n")
                markdown_lines.append(f"|{'|'.join(['---'] * len(headers))}|\n")
                for row in table.rows[1:]:
                    row_data = [cell.text for cell in row.cells]
                    markdown_lines.append(f"| {' | '.join(row_data)} |\n")
                markdown_lines.append("\n")

        markdown_lines.append("---\n\n")

    return "".join(markdown_lines)


def getfilecontent(filepath):
    file_content = ""
    try:
        ## pdf to text/markdown
        ## no wav to text/markdown as sst on client.(onnx is not cross plateform)
        if is_pdf_file(filepath):
            logging.info("{filepath} is a pdf file")
            file_content = extract_text(filepath)
            return file_content
        if is_pptx_file(filepath):
            logging.info("{filepath} is a pptx file")
            file_content = pptx_to_markdown_advanced(filepath)
            return file_content
        # default option
        logging.info("try {filepath} as text")
        with open(filepath, "r", encoding="utf-8") as f:
            file_content = f.read()
        return file_content
    except Exception:
        return file_content
